"""
Build a 10-slide visual presentation for CS516 Smart SMS.
Heavy on real screenshots; minimal text per slide.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

ROOT   = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS  = os.path.join(ROOT, "Screenshots")
OUT    = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CS516_Presentation.pptx")

# Brand palette
NAVY      = RGBColor(0x1F, 0x3A, 0x5F)
NAVY_DK   = RGBColor(0x14, 0x28, 0x45)
ACCENT    = RGBColor(0x1F, 0x6F, 0xEB)
GREEN     = RGBColor(0x2E, 0xA0, 0x43)
RED       = RGBColor(0xCF, 0x22, 0x2E)
GOLD      = RGBColor(0xE3, 0xB3, 0x41)
LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xF8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0x5F, 0x6B, 0x7C)
DARK_TEXT = RGBColor(0x1B, 0x21, 0x28)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]  # blank layout


# ---------- helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_TEXT,
             align="left", anchor="top", font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                           "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                        "right": PP_ALIGN.RIGHT}[align]
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_image(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path):
        return None
    if w is None and h is None:
        return slide.shapes.add_picture(path, x, y)
    return slide.shapes.add_picture(path, x, y, width=w, height=h)


def slide_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def slide_header(slide, title, subtitle=None):
    """Standard slide header bar"""
    add_rect(slide, 0, 0, SW, Inches(0.95), NAVY)
    add_rect(slide, 0, Inches(0.95), SW, Inches(0.04), GOLD)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.6),
             title, size=28, bold=True, color=WHITE, anchor="middle")
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.55), Inches(10), Inches(0.4),
                 subtitle, size=12, color=RGBColor(0x9B, 0xD1, 0xFF), italic=True)
    # page indicator
    add_text(slide, Inches(11.8), Inches(0.25), Inches(1.4), Inches(0.5),
             "CS516 · IAU", size=11, color=WHITE, align="right", anchor="middle")


def pill(slide, x, y, w, h, text, fill=ACCENT, color=WHITE, size=12):
    """Rounded chip"""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = color; r.font.name = "Calibri"


def feature_card(slide, x, y, w, h, icon, title, body, accent=ACCENT):
    """Feature card with colored top bar + icon + text"""
    add_rect(slide, x, y, w, h, LIGHT_BG)
    add_rect(slide, x, y, w, Inches(0.18), accent)
    # icon
    add_text(slide, x + Inches(0.15), y + Inches(0.25), Inches(0.7), Inches(0.7),
             icon, size=32, anchor="middle", align="center")
    # title
    add_text(slide, x + Inches(0.85), y + Inches(0.3), w - Inches(1.0), Inches(0.5),
             title, size=14, bold=True, color=NAVY)
    # body
    add_text(slide, x + Inches(0.2), y + Inches(0.95), w - Inches(0.4), h - Inches(1.1),
             body, size=10, color=GREY)


# =============================================================
#                    SLIDE 1 — TITLE
# =============================================================
s = prs.slides.add_slide(BLANK)
# full navy bg
add_rect(s, 0, 0, SW, SH, NAVY)
# diagonal accent stripe
stripe = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(-2), Inches(5.6),
                            Inches(20), Inches(2.5))
stripe.fill.solid(); stripe.fill.fore_color.rgb = NAVY_DK
stripe.line.fill.background()
# gold accent line
add_rect(s, Inches(0.5), Inches(2.2), Inches(0.8), Inches(0.08), GOLD)

# university line
add_text(s, Inches(0.5), Inches(0.5), Inches(8), Inches(0.4),
         "IMAM ABDULRAHMAN BIN FAISAL UNIVERSITY",
         size=12, bold=True, color=GOLD)
add_text(s, Inches(0.5), Inches(0.85), Inches(8), Inches(0.4),
         "College of Computer Science and Information Technology",
         size=11, color=RGBColor(0xB8, 0xC8, 0xDC))

# course chip
pill(s, Inches(0.5), Inches(1.5), Inches(2.4), Inches(0.45),
     "CS 516 · Term Project", fill=ACCENT)

# title
add_text(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.5),
         "Smart Student Management",
         size=54, bold=True, color=WHITE)
add_text(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.5),
         "& Attendance System",
         size=54, bold=True, color=GOLD)
add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.6),
         "A client–server application built twice — in Python and in C++",
         size=18, italic=True, color=RGBColor(0xB8, 0xC8, 0xDC))

# team strip on bottom
add_text(s, Inches(0.5), Inches(6.0), Inches(6), Inches(0.4),
         "TEAM", size=10, bold=True, color=GOLD)
add_text(s, Inches(0.5), Inches(6.35), Inches(8), Inches(0.5),
         "Ahmed Mabkhot Al-Awlaqi  ·  Abdulaziz Alghamdi  ·  Abdullah Alsaedi",
         size=14, bold=True, color=WHITE)
add_text(s, Inches(0.5), Inches(6.75), Inches(8), Inches(0.4),
         "2220007240   ·   2220003178   ·   2220002306",
         size=11, color=RGBColor(0xB8, 0xC8, 0xDC))

add_text(s, Inches(9.5), Inches(6.35), Inches(3.5), Inches(0.5),
         "Submitted to", size=10, color=GOLD, align="right")
add_text(s, Inches(9.5), Inches(6.65), Inches(3.5), Inches(0.5),
         "Mrs. Sarah Alissa", size=14, bold=True, color=WHITE, align="right")


# =============================================================
#                    SLIDE 2 — THE PROBLEM
# =============================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "The Problem", "Why we built this")

# big visual quote
add_rect(s, Inches(0.6), Inches(1.4), Inches(7.5), Inches(5.5), LIGHT_BG)
add_text(s, Inches(0.9), Inches(1.6), Inches(0.5), Inches(0.5),
         "“", size=72, bold=True, color=NAVY)
add_text(s, Inches(1.4), Inches(1.9), Inches(6.5), Inches(2.4),
         "Most departments still juggle paper attendance sheets, Excel files, "
         "and the occasional WhatsApp group. Sheets get lost. Two instructors "
         "mark the same student twice. And at the end of the semester, somebody "
         "has to count rows manually.",
         size=18, italic=True, color=DARK_TEXT)
add_text(s, Inches(1.4), Inches(4.5), Inches(6.5), Inches(0.5),
         "— our starting point",
         size=12, color=GREY)

# pain-point cards on right
add_text(s, Inches(8.5), Inches(1.4), Inches(4.5), Inches(0.5),
         "WHAT GOES WRONG TODAY", size=12, bold=True, color=NAVY)

pains = [
    ("📋", "Lost sheets", "Paper attendance can disappear before grades close."),
    ("🔁", "Duplicates", "Two instructors mark the same class twice."),
    ("⏳", "Slow reports", "End-of-semester totals are counted by hand."),
    ("📂", "No backup", "Excel files live on one laptop with no copies."),
]
y = Inches(2.0)
for icon, t, body in pains:
    add_rect(s, Inches(8.5), y, Inches(4.5), Inches(1.1), WHITE, line=RGBColor(0xE0, 0xE5, 0xEC))
    add_text(s, Inches(8.65), y + Inches(0.15), Inches(0.6), Inches(0.8),
             icon, size=24, anchor="middle")
    add_text(s, Inches(9.3), y + Inches(0.1), Inches(3.0), Inches(0.4),
             t, size=14, bold=True, color=NAVY)
    add_text(s, Inches(9.3), y + Inches(0.45), Inches(3.0), Inches(0.6),
             body, size=10, color=GREY)
    y += Inches(1.25)


# =============================================================
#                    SLIDE 3 — WHAT WE BUILT
# =============================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "What We Built", "The same system, in two languages")

# big subhead
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.6),
         "One server holds the data. Two clients let you work with it.",
         size=18, color=NAVY, bold=True)

# two halves with screenshots
LEFT_X = Inches(0.5)
RIGHT_X = Inches(6.95)
TOP_Y = Inches(2.0)
CARD_W = Inches(6.0)
CARD_H = Inches(5.0)

# Python card
add_rect(s, LEFT_X, TOP_Y, CARD_W, CARD_H, LIGHT_BG)
add_rect(s, LEFT_X, TOP_Y, CARD_W, Inches(0.45), ACCENT)
add_text(s, LEFT_X + Inches(0.2), TOP_Y + Inches(0.08), Inches(5), Inches(0.4),
         "PYTHON  ·  Tkinter GUI + SQLite", size=14, bold=True, color=WHITE)
add_image(s, os.path.join(SHOTS, "py_02_main_students.png"),
          LEFT_X + Inches(0.2), TOP_Y + Inches(0.7),
          w=Inches(5.6))

# C++ card
add_rect(s, RIGHT_X, TOP_Y, CARD_W, CARD_H, LIGHT_BG)
add_rect(s, RIGHT_X, TOP_Y, CARD_W, Inches(0.45), GREEN)
add_text(s, RIGHT_X + Inches(0.2), TOP_Y + Inches(0.08), Inches(5), Inches(0.4),
         "C++  ·  OOP console + flat files", size=14, bold=True, color=WHITE)
add_image(s, os.path.join(SHOTS, "cpp_03_students_view_all.png"),
          RIGHT_X + Inches(0.2), TOP_Y + Inches(0.7),
          w=Inches(5.6))

# bottom note
add_text(s, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
         "Both clients speak the same protocol — they could even share the server.",
         size=11, italic=True, color=GREY, align="center")


# =============================================================
#                    SLIDE 4 — FEATURES
# =============================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Features", "Everything Smart SMS does")

# left: feature grid
features = [
    ("🔐", "Role-based Login", "Admin sees everything; instructor sees what they teach.", ACCENT),
    ("👤", "Students CRUD", "Add, update, delete, search — full lifecycle.", GREEN),
    ("📚", "Courses CRUD", "Same surface for courses with code as primary key.", GOLD),
    ("📝", "Enrollment", "Per-semester enrollment with duplicate protection.", ACCENT),
    ("✅", "Attendance", "Mark Present/Absent, validated against enrollment.", GREEN),
    ("📊", "Summary Reports", "Per-course present/absent counts in one click.", RED),
]

CARD_W = Inches(2.6)
CARD_H = Inches(2.0)
ROW1 = Inches(1.4)
ROW2 = Inches(3.55)
COL_X = [Inches(0.5), Inches(3.25), Inches(6.0)]

for i, (icon, t, body, accent) in enumerate(features):
    col = i % 3
    row = i // 3
    x = COL_X[col]
    y = ROW1 if row == 0 else ROW2
    feature_card(s, x, y, CARD_W, CARD_H, icon, t, body, accent=accent)

# right: attendance summary as the "money shot"
add_text(s, Inches(8.9), Inches(1.4), Inches(4.2), Inches(0.4),
         "WHAT INSTRUCTORS CARE ABOUT", size=11, bold=True, color=NAVY)
add_image(s, os.path.join(SHOTS, "py_10_attendance_summary.png"),
          Inches(8.9), Inches(1.85), w=Inches(4.2))
add_text(s, Inches(8.9), Inches(5.6), Inches(4.2), Inches(0.5),
         "Per-course summary", size=14, bold=True, color=NAVY)
add_text(s, Inches(8.9), Inches(5.95), Inches(4.2), Inches(1.1),
         "Counts present vs absent for every student in a course. "
         "What used to take 10 minutes of counting now takes one click.",
         size=10, color=GREY)


# =============================================================
#                    SLIDE 5 — ARCHITECTURE
# =============================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Architecture", "One server, two clients, one protocol")

# Server box (center)
SX = Inches(5.3); SY = Inches(2.3); SW2 = Inches(2.7); SH2 = Inches(2.4)
add_rect(s, SX, SY, SW2, SH2, NAVY)
add_text(s, SX, SY + Inches(0.3), SW2, Inches(0.6),
         "SERVER", size=18, bold=True, color=WHITE, align="center")
add_text(s, SX, SY + Inches(0.9), SW2, Inches(0.6),
         "🖥️", size=36, align="center", anchor="middle")
add_text(s, SX, SY + Inches(1.55), SW2, Inches(0.4),
         "Validation", size=11, color=GOLD, align="center", bold=True)
add_text(s, SX, SY + Inches(1.85), SW2, Inches(0.4),
         "Storage", size=11, color=GOLD, align="center", bold=True)

# Python client (left)
add_rect(s, Inches(0.6), SY, Inches(3.2), SH2, LIGHT_BG)
add_rect(s, Inches(0.6), SY, Inches(3.2), Inches(0.4), ACCENT)
add_text(s, Inches(0.6), SY + Inches(0.05), Inches(3.2), Inches(0.35),
         "PYTHON CLIENT", size=12, bold=True, color=WHITE, align="center")
add_text(s, Inches(0.6), SY + Inches(0.55), Inches(3.2), Inches(0.5),
         "🪟", size=36, align="center")
add_text(s, Inches(0.6), SY + Inches(1.15), Inches(3.2), Inches(0.4),
         "Tkinter GUI", size=13, bold=True, color=NAVY, align="center")
add_text(s, Inches(0.6), SY + Inches(1.5), Inches(3.2), Inches(0.4),
         "JSON over TCP", size=10, color=GREY, align="center", italic=True)
add_text(s, Inches(0.6), SY + Inches(1.85), Inches(3.2), Inches(0.4),
         "127.0.0.1 : 5055", size=10, color=GREY, align="center")

# C++ client (right)
add_rect(s, Inches(9.5), SY, Inches(3.2), SH2, LIGHT_BG)
add_rect(s, Inches(9.5), SY, Inches(3.2), Inches(0.4), GREEN)
add_text(s, Inches(9.5), SY + Inches(0.05), Inches(3.2), Inches(0.35),
         "C++ CLIENT", size=12, bold=True, color=WHITE, align="center")
add_text(s, Inches(9.5), SY + Inches(0.55), Inches(3.2), Inches(0.5),
         "💻", size=36, align="center")
add_text(s, Inches(9.5), SY + Inches(1.15), Inches(3.2), Inches(0.4),
         "Console UI", size=13, bold=True, color=NAVY, align="center")
add_text(s, Inches(9.5), SY + Inches(1.5), Inches(3.2), Inches(0.4),
         "ACTION|k=v over TCP", size=10, color=GREY, align="center", italic=True)
add_text(s, Inches(9.5), SY + Inches(1.85), Inches(3.2), Inches(0.4),
         "127.0.0.1 : 5066", size=10, color=GREY, align="center")

# arrows
arrow1 = s.shapes.add_connector(1, Inches(3.8), SY + Inches(1.2), Inches(5.3), SY + Inches(1.2))
arrow1.line.color.rgb = NAVY; arrow1.line.width = Pt(3)
# add arrowhead via XML
ln = arrow1.line._get_or_add_ln()
tail = etree.SubElement(ln, qn("a:tailEnd")); tail.set("type", "triangle"); tail.set("w", "lg"); tail.set("len", "lg")

arrow2 = s.shapes.add_connector(1, Inches(8.0), SY + Inches(1.2), Inches(9.5), SY + Inches(1.2))
arrow2.line.color.rgb = NAVY; arrow2.line.width = Pt(3)
ln = arrow2.line._get_or_add_ln()
tail = etree.SubElement(ln, qn("a:tailEnd")); tail.set("type", "triangle"); tail.set("w", "lg"); tail.set("len", "lg")
head = etree.SubElement(ln, qn("a:headEnd")); head.set("type", "triangle"); head.set("w", "lg"); head.set("len", "lg")

# Storage layer at bottom
add_rect(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.7), LIGHT_BG)
add_text(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.4),
         "STORAGE  ·  same data model, two engines",
         size=12, bold=True, color=NAVY, align="center")
# Two storage cards
add_rect(s, Inches(1.5), Inches(6.05), Inches(4.5), Inches(1.0), WHITE, line=ACCENT)
add_text(s, Inches(1.5), Inches(6.15), Inches(4.5), Inches(0.4),
         "🐍  SQLite (mgmt via sqlite3)", size=13, bold=True, color=NAVY, align="center")
add_text(s, Inches(1.5), Inches(6.55), Inches(4.5), Inches(0.4),
         "Parameterised SQL, JOINs, GROUP BY", size=10, color=GREY, align="center", italic=True)

add_rect(s, Inches(7.3), Inches(6.05), Inches(4.5), Inches(1.0), WHITE, line=GREEN)
add_text(s, Inches(7.3), Inches(6.15), Inches(4.5), Inches(0.4),
         "⚙️  TAB-separated flat files", size=13, bold=True, color=NAVY, align="center")
add_text(s, Inches(7.3), Inches(6.55), Inches(4.5), Inches(0.4),
         "One file per entity, std::vector in memory", size=10, color=GREY, align="center", italic=True)


# =============================================================
#                    SLIDE 6 — PYTHON DEEP DIVE
# =============================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Python Implementation", "Tkinter GUI · SQLite · TCP server (multi-threaded)")

# Three screenshot strip
add_image(s, os.path.join(SHOTS, "py_01_login.png"),
          Inches(0.5), Inches(1.35), w=Inches(3.2))
add_image(s, os.path.join(SHOTS, "py_08_enrollment_tab.png"),
          Inches(3.85), Inches(1.35), w=Inches(4.6))
add_image(s, os.path.join(SHOTS, "py_09_attendance_report.png"),
          Inches(8.55), Inches(1.35), w=Inches(4.4))

# Caption row
for x, lbl in [(Inches(0.5), "1. Sign in"),
               (Inches(3.85), "2. Enroll students"),
               (Inches(8.55), "3. Track attendance")]:
    pill(s, x, Inches(4.4), Inches(2.0), Inches(0.35), lbl, fill=ACCENT, size=10)

# Tech stack chips
add_text(s, Inches(0.5), Inches(5.05), Inches(4), Inches(0.4),
         "TECH STACK", size=11, bold=True, color=NAVY)
chips = [("tkinter", ACCENT), ("sqlite3", ACCENT), ("socket", ACCENT),
         ("threading", ACCENT), ("json", ACCENT)]
x = Inches(0.5)
for t, c in chips:
    pill(s, x, Inches(5.45), Inches(1.3), Inches(0.4), t, fill=c, size=11)
    x += Inches(1.4)

# Highlights right
add_text(s, Inches(8.0), Inches(5.05), Inches(5), Inches(0.4),
         "WHY PYTHON SHINES HERE", size=11, bold=True, color=NAVY)
highlights = [
    "✓ Zero external dependencies — standard library only",
    "✓ GUI built in ~600 lines of tkinter",
    "✓ SQL handles GROUP BY for free",
    "✓ Cross-platform from day one",
]
y = Inches(5.45)
for h in highlights:
    add_text(s, Inches(8.0), y, Inches(5.3), Inches(0.35),
             h, size=11, color=DARK_TEXT)
    y += Inches(0.35)


# =============================================================
#                    SLIDE 7 — C++ DEEP DIVE
# =============================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "C++ Implementation", "OOP server · Winsock · static .exe — zero dependencies")

# Three screenshot strip
add_image(s, os.path.join(SHOTS, "cpp_01_login.png"),
          Inches(0.5), Inches(1.35), w=Inches(4.0))
add_image(s, os.path.join(SHOTS, "cpp_03_students_view_all.png"),
          Inches(4.65), Inches(1.35), w=Inches(4.2))
add_image(s, os.path.join(SHOTS, "cpp_10_attendance_summary.png"),
          Inches(9.0), Inches(1.35), w=Inches(3.95))

for x, lbl in [(Inches(0.5), "1. Login"),
               (Inches(4.65), "2. View students"),
               (Inches(9.0), "3. Summary report")]:
    pill(s, x, Inches(4.4), Inches(2.0), Inches(0.35), lbl, fill=GREEN, size=10)

# Tech stack
add_text(s, Inches(0.5), Inches(5.05), Inches(4), Inches(0.4),
         "TECH STACK", size=11, bold=True, color=NAVY)
chips = [("C++17", GREEN), ("Winsock2", GREEN), ("std::thread", GREEN),
         ("std::filesystem", GREEN), ("MinGW g++ 14.2", GREEN)]
x = Inches(0.5)
for t, c in chips:
    pill(s, x, Inches(5.45), Inches(1.7), Inches(0.4), t, fill=c, size=10)
    x += Inches(1.8)

# Highlights right
add_text(s, Inches(8.0), Inches(5.05), Inches(5), Inches(0.4),
         "WHY C++ SHINES HERE", size=11, bold=True, color=NAVY)
highlights = [
    "✓ Single self-contained .exe (no DLLs needed)",
    "✓ Clean OOP: one Repository class per entity",
    "✓ Predictable memory via RAII — no GC pauses",
    "✓ Static type checking catches bugs at compile time",
]
y = Inches(5.45)
for h in highlights:
    add_text(s, Inches(8.0), y, Inches(5.3), Inches(0.35),
             h, size=11, color=DARK_TEXT)
    y += Inches(0.35)


# =============================================================
#                    SLIDE 8 — COMPARISON PART 1
# =============================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Comparison — Part 1", "Syntax · Dependencies · Paradigm · Memory")

rows = [
    ("Syntax",
     "Indented blocks, no braces.\nDynamic typing.\nprint(\"hi\")",
     "Braces and semicolons.\nStatic typing.\nstd::cout << \"hi\";"),
    ("Dependencies",
     "Standard library only.\ntkinter, sqlite3, socket all bundled.",
     "Requires C++17 compiler\n+ ws2_32 (Winsock) on Windows."),
    ("Paradigm",
     "Multi-paradigm: procedural,\nOOP, and functional all welcome.",
     "OOP + procedural primary;\ngeneric programming via templates."),
    ("Memory",
     "Automatic — reference counting\n+ cyclic GC. Developer never frees.",
     "Manual but rescued by RAII\n(std::string, std::vector free in scope)."),
]

# Header row
HX = Inches(0.6); HY = Inches(1.35)
COL_W = [Inches(2.3), Inches(4.85), Inches(4.85)]
TOTAL_W = sum(COL_W, Emu(0))
HEAD_H = Inches(0.55)
ROW_H = Inches(1.2)

# header
xacc = HX
for i, (txt, w) in enumerate(zip(["Feature", "🐍 Python", "⚙️ C++"], COL_W)):
    fill = NAVY if i == 0 else (ACCENT if i == 1 else GREEN)
    add_rect(s, xacc, HY, w, HEAD_H, fill)
    add_text(s, xacc, HY + Inches(0.1), w, Inches(0.4),
             txt, size=15, bold=True, color=WHITE, align="center")
    xacc += w

# rows
y = HY + HEAD_H
for i, (feat, py, cpp) in enumerate(rows):
    fill = LIGHT_BG if i % 2 == 0 else WHITE
    xacc = HX
    for j, (txt, w) in enumerate(zip([feat, py, cpp], COL_W)):
        add_rect(s, xacc, y, w, ROW_H, fill, line=RGBColor(0xE0, 0xE5, 0xEC))
        is_feat = (j == 0)
        add_text(s, xacc + Inches(0.15), y + Inches(0.1), w - Inches(0.3), ROW_H - Inches(0.2),
                 txt, size=(13 if is_feat else 11),
                 bold=is_feat, color=(NAVY if is_feat else DARK_TEXT),
                 anchor="middle")
        xacc += w
    y += ROW_H


# =============================================================
#               SLIDE 9 — COMPARISON PART 2 (CODE)
# =============================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Comparison — Part 2", "Same function, two languages")

# Section: Add Student
add_text(s, Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.4),
         "Function: Add a Student   (validation + persistence)",
         size=15, bold=True, color=NAVY)

# Python code block (left)
add_rect(s, Inches(0.6), Inches(1.75), Inches(6.05), Inches(2.5), RGBColor(0x12, 0x1B, 0x2D))
add_rect(s, Inches(0.6), Inches(1.75), Inches(6.05), Inches(0.32), ACCENT)
add_text(s, Inches(0.7), Inches(1.78), Inches(5), Inches(0.3),
         "🐍 Python  ·  server.py", size=10, bold=True, color=WHITE)
py_code = ("if not is_valid_id(p['student_id']):\n"
           "    return {'ok': False,\n"
           "            'error': 'Student ID must be 10 digits'}\n"
           "conn.execute(\n"
           "  \"INSERT INTO students(...) VALUES(?,?,?,?,?,?)\",\n"
           "  (sid, name, email, phone, major, level))\n"
           "conn.commit()\n"
           "return {'ok': True}")
add_text(s, Inches(0.75), Inches(2.15), Inches(5.8), Inches(2.0),
         py_code, size=11, color=RGBColor(0xE0, 0xE8, 0xF4), font="Consolas")

# C++ code block (right)
add_rect(s, Inches(6.85), Inches(1.75), Inches(6.05), Inches(2.5), RGBColor(0x12, 0x1B, 0x2D))
add_rect(s, Inches(6.85), Inches(1.75), Inches(6.05), Inches(0.32), GREEN)
add_text(s, Inches(6.95), Inches(1.78), Inches(5), Inches(0.3),
         "⚙️ C++  ·  server.cpp", size=10, bold=True, color=WHITE)
cpp_code = ("if (!validId(sid))\n"
            "    return err(\"Student ID must be 10 digits\");\n"
            "auto e = students_.add(\n"
            "    {sid, name, email, phone, major, level});\n"
            "if (!e.empty()) return err(e);\n"
            "return ok();\n"
            "// repository writes a TAB-separated row\n"
            "// to data/students.txt")
add_text(s, Inches(7.00), Inches(2.15), Inches(5.8), Inches(2.0),
         cpp_code, size=11, color=RGBColor(0xE0, 0xE8, 0xF4), font="Consolas")

# Function: Attendance summary
add_text(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(0.4),
         "Function: Per-course Attendance Summary",
         size=15, bold=True, color=NAVY)

# Python
add_rect(s, Inches(0.6), Inches(4.9), Inches(6.05), Inches(2.2), RGBColor(0x12, 0x1B, 0x2D))
add_rect(s, Inches(0.6), Inches(4.9), Inches(6.05), Inches(0.32), ACCENT)
add_text(s, Inches(0.7), Inches(4.93), Inches(5), Inches(0.3),
         "🐍 Python  ·  one SQL query", size=10, bold=True, color=WHITE)
py2 = ("SELECT student_id, name,\n"
       "  SUM(status='Present') AS present_count,\n"
       "  SUM(status='Absent')  AS absent_count\n"
       "FROM attendance JOIN students USING(student_id)\n"
       "WHERE course_code = ?\n"
       "GROUP BY student_id;")
add_text(s, Inches(0.75), Inches(5.3), Inches(5.8), Inches(1.75),
         py2, size=11, color=RGBColor(0xE0, 0xE8, 0xF4), font="Consolas")

# C++
add_rect(s, Inches(6.85), Inches(4.9), Inches(6.05), Inches(2.2), RGBColor(0x12, 0x1B, 0x2D))
add_rect(s, Inches(6.85), Inches(4.9), Inches(6.05), Inches(0.32), GREEN)
add_text(s, Inches(6.95), Inches(4.93), Inches(5), Inches(0.3),
         "⚙️ C++  ·  manual aggregation", size=10, bold=True, color=WHITE)
cpp2 = ("map<string, pair<int,int>> counts;\n"
        "for (auto& row : attendance_.all()) {\n"
        "  if (row[1] != code) continue;\n"
        "  if (row[3] == \"Present\") counts[row[0]].first++;\n"
        "  else                       counts[row[0]].second++;\n"
        "}")
add_text(s, Inches(7.00), Inches(5.3), Inches(5.8), Inches(1.75),
         cpp2, size=11, color=RGBColor(0xE0, 0xE8, 0xF4), font="Consolas")

# Verdict pill at bottom
pill(s, Inches(3.0), Inches(7.05), Inches(7.3), Inches(0.4),
     "Python wins on conciseness · C++ wins on transparency", fill=NAVY, size=12)


# =============================================================
#               SLIDE 10 — VERDICT + DEMO + Q&A
# =============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(5.6), SW, Inches(2), NAVY_DK)

# Big title
add_text(s, Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.8),
         "Wrapping Up", size=44, bold=True, color=WHITE)
add_rect(s, Inches(0.5), Inches(1.4), Inches(1.0), Inches(0.08), GOLD)

# Verdict block
add_text(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.5),
         "THE VERDICT", size=11, bold=True, color=GOLD)

verdict = (
    "Python was faster to write and the GUI looked polished with very little effort. "
    "C++ forced us to be explicit and produced a smaller, dependency-free binary, "
    "but every feature took noticeably longer to land. They are not in competition — "
    "they answer different questions."
)
add_text(s, Inches(0.5), Inches(2.05), Inches(12.3), Inches(2.0),
         verdict, size=18, color=WHITE, italic=True)

# Stats row
stats = [
    ("2", "languages", GOLD),
    ("17", "server actions", ACCENT),
    ("21", "screenshots", GREEN),
    ("34", "files in repo", RED),
]
x = Inches(0.5); w = Inches(3.0); gap = Inches(0.13)
for val, label, col in stats:
    add_rect(s, x, Inches(4.1), w, Inches(1.2), NAVY_DK)
    add_rect(s, x, Inches(4.1), Inches(0.08), Inches(1.2), col)
    add_text(s, x + Inches(0.2), Inches(4.2), w - Inches(0.3), Inches(0.7),
             val, size=36, bold=True, color=col, anchor="middle")
    add_text(s, x + Inches(0.2), Inches(4.85), w - Inches(0.3), Inches(0.35),
             label.upper(), size=11, color=WHITE, anchor="middle")
    x += w + gap

# Bottom strip
add_text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
         "Live Demo  ·  Q & A", size=28, bold=True, color=GOLD, align="center")
add_text(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.4),
         "github.com/abdullahalsaedi2003-dotcom/adaved-programing-project",
         size=12, color=RGBColor(0xB8, 0xC8, 0xDC), align="center", italic=True)
add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
         "Thank you  ·  Ahmed · Abdulaziz · Abdullah",
         size=14, bold=True, color=WHITE, align="center")


# ============== save ==============
prs.save(OUT)
print("Saved:", OUT, "(", os.path.getsize(OUT), "bytes )")
